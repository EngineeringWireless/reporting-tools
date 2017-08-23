from pptx import Presentation
from pptx.util import Length
from utils.prs_tools import get_layout
from utils.image_tools import fit_image


class Report():

    '''
    creates a class containing a pptx presentation used to create the report.
    '''
    def __init__(self):

        self.prs = Presentation('tools/PPTX_TEMPLATE.pptx')

    '''
    adds a "Plot with Legend" slide to the presentationa and returns that slide
    to be worked with.
    '''
    def add_image_slide(self):
        img_layout = 'Plot with Legend'
        self.cur_slide = self.prs.slides.add_slide(
                                get_layout(self.prs.slide_layouts, img_layout)
                                )
        return self.cur_slide

    def insert_image(self, placeholder, image):

        size = [
            int(Length(placeholder.width).pt),
            int(Length(placeholder.height).pt)
        ]

        img = fit_image(image, size)

        img_placeholder = placeholder.insert_picture(img)

        return img_placeholder

    '''
    saves presentation to given location.
    '''
    def save(self, file_name):
        self.prs.save(file_name)
